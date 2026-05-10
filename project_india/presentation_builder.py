from __future__ import annotations

import json
from pathlib import Path
import re

from project_india import paths
from project_india.topics import TOPIC_FOLDERS, create_topic, slugify


def _read_section(markdown: str, heading: str) -> str:
    marker = f"## {heading}"
    if marker not in markdown:
        return "TBD."
    after = markdown.split(marker, 1)[1]
    next_heading = after.find("\n## ")
    section = after if next_heading == -1 else after[:next_heading]
    lines = [line.strip() for line in section.splitlines() if line.strip()]
    return "\n".join(lines) if lines else "TBD."


def _bullets(text: str, limit: int = 5) -> list[str]:
    bullets = []
    for line in text.splitlines():
        clean = line.strip().lstrip("-").strip()
        if clean and clean != "TBD.":
            bullets.append(clean)
    if not bullets and text.strip() and text.strip() != "TBD.":
        bullets = [text.strip()]
    return bullets[:limit] or ["TBD."]


def _topic_path(slug: str, category: str) -> Path:
    return TOPIC_FOLDERS[category] / f"{slug}.md"


def _load_or_create_topic(title: str, slug: str, category: str) -> Path:
    path = _topic_path(slug, category)
    if not path.exists():
        create_topic(title, slug=slug, category=category)
    return path


def _outline_path(slug: str) -> Path:
    return paths.PRESENTATIONS / f"{slug}-outline.md"


def _topic_data_path(slug: str) -> Path:
    return paths.TOPIC_DATA / f"{slug}.json"


def _load_topic_data(slug: str) -> dict:
    path = _topic_data_path(slug)
    if not path.exists():
        return {}
    try:
        return json.loads(path.read_text(encoding="utf-8"))
    except json.JSONDecodeError:
        return {}


def _number(value: object) -> float | None:
    if isinstance(value, int | float):
        return float(value)
    if isinstance(value, str):
        cleaned = value.replace(",", "").strip()
        try:
            return float(cleaned)
        except ValueError:
            return None
    return None


def _outline_slides(markdown: str) -> list[tuple[str, list[str]]]:
    outline = _read_section(markdown, "Slide Outline")
    if outline == "TBD.":
        return []

    slides: list[tuple[str, list[str]]] = []
    current_title: str | None = None
    current_bullets: list[str] = []

    for line in outline.splitlines():
        stripped = line.strip()
        match = re.match(r"^\d+\.\s+(.*)", stripped)
        if match:
            if current_title:
                slides.append((current_title, current_bullets[:5] or ["TBD."]))
            current_title = match.group(1).strip()
            current_bullets = []
            continue

        if current_title and stripped.startswith("-"):
            current_bullets.append(stripped.lstrip("-").strip())
        elif current_title and stripped:
            current_bullets.append(stripped)

    if current_title:
        slides.append((current_title, current_bullets[:5] or ["TBD."]))

    return slides[:12]


def build_presentation(
    title: str,
    slug: str | None = None,
    category: str = "sectors",
    output_path: Path | None = None,
) -> Path:
    topic_slug = slugify(slug or title)
    topic_path = _load_or_create_topic(title, topic_slug, category)
    markdown = topic_path.read_text(encoding="utf-8")
    outline_path = _outline_path(topic_slug)
    outline_markdown = outline_path.read_text(encoding="utf-8") if outline_path.exists() else ""
    topic_data = _load_topic_data(topic_slug)

    try:
        from pptx import Presentation
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
        from pptx.util import Inches, Pt
    except ImportError as error:
        raise SystemExit(
            "python-pptx is required. Install with: python3 -m pip install -e '.[presentation]'"
        ) from error

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    ink = RGBColor(21, 21, 21)
    muted = RGBColor(103, 98, 90)
    saffron = RGBColor(217, 119, 6)
    paper = RGBColor(247, 241, 231)

    def set_background(slide) -> None:
        bg = slide.background.fill
        bg.solid()
        bg.fore_color.rgb = paper

    def add_footer(slide, note: str = "Project India | generated presentation draft | verify time-sensitive claims") -> None:
        footer = slide.shapes.add_textbox(Inches(0.65), Inches(6.95), Inches(11.7), Inches(0.22))
        footer_tf = footer.text_frame
        footer_tf.text = note
        footer_para = footer_tf.paragraphs[0]
        footer_para.alignment = PP_ALIGN.LEFT
        footer_run = footer_para.runs[0]
        footer_run.font.size = Pt(8)
        footer_run.font.color.rgb = muted

    def add_header(slide, kicker: str, claim: str) -> None:
        marker = slide.shapes.add_shape(1, Inches(0.65), Inches(0.55), Inches(0.18), Inches(0.18))
        marker.fill.solid()
        marker.fill.fore_color.rgb = saffron
        marker.line.fill.background()

        kicker_box = slide.shapes.add_textbox(Inches(0.95), Inches(0.48), Inches(2.2), Inches(0.3))
        kicker_tf = kicker_box.text_frame
        kicker_tf.text = kicker.upper()
        kicker_run = kicker_tf.paragraphs[0].runs[0]
        kicker_run.font.size = Pt(9)
        kicker_run.font.bold = True
        kicker_run.font.color.rgb = muted

        title_box = slide.shapes.add_textbox(Inches(0.65), Inches(1.02), Inches(11.3), Inches(1.0))
        title_tf = title_box.text_frame
        title_tf.word_wrap = True
        title_tf.text = claim
        title_run = title_tf.paragraphs[0].runs[0]
        title_run.font.size = Pt(27)
        title_run.font.bold = True
        title_run.font.color.rgb = ink

    def add_slide(kicker: str, claim: str, bullets: list[str]) -> None:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        set_background(slide)
        add_header(slide, kicker, claim)

        body_box = slide.shapes.add_textbox(Inches(1.05), Inches(2.45), Inches(10.7), Inches(3.5))
        body_tf = body_box.text_frame
        body_tf.word_wrap = True
        body_tf.clear()
        for index, item in enumerate(bullets):
            para = body_tf.paragraphs[0] if index == 0 else body_tf.add_paragraph()
            para.text = item
            para.level = 0
            para.space_after = Pt(10)
            para.font.size = Pt(18)
            para.font.color.rgb = ink

        add_footer(slide)

    def add_metric_slide(metrics: list[dict]) -> None:
        if not metrics:
            return
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        set_background(slide)
        add_header(slide, "Evidence", "The core evidence should be visible before the argument gets abstract.")
        for idx, metric in enumerate(metrics[:6]):
            x = Inches(0.8 + (idx % 3) * 4.1)
            y = Inches(2.25 + (idx // 3) * 1.55)
            box = slide.shapes.add_shape(1, x, y, Inches(3.55), Inches(1.16))
            box.fill.solid()
            box.fill.fore_color.rgb = RGBColor(255, 250, 240)
            box.line.color.rgb = RGBColor(210, 196, 174)
            value = metric.get("value", "TBD")
            unit = metric.get("unit", "")
            label = metric.get("label", "Metric")
            context = metric.get("context") or metric.get("date") or ""
            value_box = slide.shapes.add_textbox(x + Inches(0.18), y + Inches(0.14), Inches(3.1), Inches(0.35))
            value_tf = value_box.text_frame
            value_tf.text = f"{value} {unit}".strip()
            value_run = value_tf.paragraphs[0].runs[0]
            value_run.font.size = Pt(22)
            value_run.font.bold = True
            value_run.font.color.rgb = saffron
            label_box = slide.shapes.add_textbox(x + Inches(0.18), y + Inches(0.55), Inches(3.1), Inches(0.28))
            label_tf = label_box.text_frame
            label_tf.text = str(label)
            label_run = label_tf.paragraphs[0].runs[0]
            label_run.font.size = Pt(11)
            label_run.font.bold = True
            label_run.font.color.rgb = ink
            context_box = slide.shapes.add_textbox(x + Inches(0.18), y + Inches(0.83), Inches(3.1), Inches(0.22))
            context_tf = context_box.text_frame
            context_tf.text = str(context)[:90]
            context_run = context_tf.paragraphs[0].runs[0]
            context_run.font.size = Pt(8.5)
            context_run.font.color.rgb = muted
        add_footer(slide, "Sources embedded in topic_data JSON; verify official/current figures before publication.")

    def add_bar_slide(comparison: dict) -> None:
        items = comparison.get("items") or []
        numeric_items = [(str(item.get("label", "")), _number(item.get("value"))) for item in items]
        numeric_items = [(label, value) for label, value in numeric_items if value is not None]
        if len(numeric_items) < 2:
            return
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        set_background(slide)
        add_header(slide, "Comparison", comparison.get("title") or "The comparison needs numbers, not just claims.")
        max_value = max(value for _, value in numeric_items) or 1
        for idx, (label, value) in enumerate(numeric_items[:7]):
            y = Inches(2.15 + idx * 0.58)
            label_box = slide.shapes.add_textbox(Inches(0.95), y, Inches(2.75), Inches(0.25))
            label_tf = label_box.text_frame
            label_tf.text = label
            label_run = label_tf.paragraphs[0].runs[0]
            label_run.font.size = Pt(12)
            label_run.font.bold = True
            label_run.font.color.rgb = ink
            bg = slide.shapes.add_shape(1, Inches(3.85), y, Inches(6.8), Inches(0.25))
            bg.fill.solid()
            bg.fill.fore_color.rgb = RGBColor(229, 219, 203)
            bg.line.fill.background()
            bar = slide.shapes.add_shape(1, Inches(3.85), y, Inches(6.8 * (value / max_value)), Inches(0.25))
            bar.fill.solid()
            bar.fill.fore_color.rgb = saffron if idx == 0 else RGBColor(49, 90, 140)
            bar.line.fill.background()
            value_box = slide.shapes.add_textbox(Inches(10.85), y - Inches(0.03), Inches(1.1), Inches(0.32))
            value_tf = value_box.text_frame
            value_tf.text = f"{value:g} {comparison.get('unit', '')}".strip()
            value_run = value_tf.paragraphs[0].runs[0]
            value_run.font.size = Pt(12)
            value_run.font.bold = True
            value_run.font.color.rgb = ink
        add_footer(slide, comparison.get("source") or "Source recorded in topic_data JSON.")

    def add_timeline_slide(events: list[dict]) -> None:
        if not events:
            return
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        set_background(slide)
        add_header(slide, "Timeline", "The story should be anchored in dates, not only interpretation.")
        usable = events[:6]
        left = 0.85
        width = 11.3
        step = width / max(1, len(usable) - 1)
        y = 3.25
        line = slide.shapes.add_shape(1, Inches(left), Inches(y), Inches(width), Inches(0.03))
        line.fill.solid()
        line.fill.fore_color.rgb = RGBColor(207, 196, 178)
        line.line.fill.background()
        for idx, event in enumerate(usable):
            x = left + idx * step
            dot = slide.shapes.add_shape(1, Inches(x), Inches(y - 0.08), Inches(0.18), Inches(0.18))
            dot.fill.solid()
            dot.fill.fore_color.rgb = saffron
            dot.line.fill.background()
            date_box = slide.shapes.add_textbox(Inches(x - 0.35), Inches(y - 0.62), Inches(0.9), Inches(0.25))
            date_tf = date_box.text_frame
            date_tf.text = str(event.get("date", "TBD"))
            date_run = date_tf.paragraphs[0].runs[0]
            date_run.font.size = Pt(10)
            date_run.font.bold = True
            date_run.font.color.rgb = ink
            event_box = slide.shapes.add_textbox(Inches(x - 0.55), Inches(y + 0.28), Inches(1.45), Inches(0.75))
            event_tf = event_box.text_frame
            event_tf.word_wrap = True
            event_tf.text = str(event.get("event", "Event"))[:90]
            event_run = event_tf.paragraphs[0].runs[0]
            event_run.font.size = Pt(9)
            event_run.font.color.rgb = muted
        add_footer(slide, "Timeline generated from structured topic_data; check source log for provenance.")

    def add_table_slide(table_data: dict) -> None:
        columns = table_data.get("columns") or []
        rows = table_data.get("rows") or []
        if not columns or not rows:
            return
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        set_background(slide)
        add_header(slide, "Data Table", table_data.get("title") or "Structured data should travel with the argument.")
        row_count = min(len(rows), 6) + 1
        col_count = min(len(columns), 4)
        table = slide.shapes.add_table(row_count, col_count, Inches(0.85), Inches(2.15), Inches(11.6), Inches(3.75)).table
        for col_idx, column in enumerate(columns[:col_count]):
            cell = table.cell(0, col_idx)
            cell.text = str(column)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(36, 33, 29)
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.color.rgb = RGBColor(255, 255, 255)
                    run.font.bold = True
                    run.font.size = Pt(10)
        for row_idx, row in enumerate(rows[: row_count - 1], start=1):
            for col_idx in range(col_count):
                cell = table.cell(row_idx, col_idx)
                value = row[col_idx] if col_idx < len(row) else ""
                cell.text = str(value)[:120]
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(255, 250, 240) if row_idx % 2 else paper
                for paragraph in cell.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.color.rgb = ink
                        run.font.size = Pt(9)
        add_footer(slide, table_data.get("source") or "Source recorded in topic_data JSON.")

    def add_gap_slide(gaps: list[str]) -> None:
        if not gaps:
            return
        add_slide("Data gaps", "The deck should be honest about what is still missing.", gaps[:6])

    researched_slides = _outline_slides(outline_markdown)
    metrics = topic_data.get("metrics") or []
    comparisons = topic_data.get("comparisons") or []
    timeline = topic_data.get("timeline") or []
    tables = topic_data.get("tables") or []
    data_gaps = topic_data.get("data_gaps") or []

    if researched_slides:
        core_message = _read_section(outline_markdown, "Core Message")
        add_slide("Topic", title, _bullets(core_message, 3))
        add_metric_slide(metrics)
        if comparisons:
            add_bar_slide(comparisons[0])
        add_timeline_slide(timeline)
        if tables:
            add_table_slide(tables[0])
        for index, (claim, bullets) in enumerate(researched_slides, start=1):
            add_slide(f"Slide {index}", claim, bullets)
        add_gap_slide(data_gaps)
    else:
        add_slide("Topic", title, [_read_section(markdown, "One-Line Summary")])
        add_metric_slide(metrics)
        if comparisons:
            add_bar_slide(comparisons[0])
        add_timeline_slide(timeline)
        if tables:
            add_table_slide(tables[0])
        add_slide("Why it matters", "This topic matters because it changes the strategic picture.", _bullets(_read_section(markdown, "Why This Matters"), 4))
        add_slide("Current state", "The current state gives us the baseline for analysis.", _bullets(_read_section(markdown, "Current State"), 5))
        add_slide("Context", "Historical context prevents shallow conclusions.", _bullets(_read_section(markdown, "Historical Context"), 5))
        add_slide("Actors", "The key actors define incentives and constraints.", _bullets(_read_section(markdown, "Key Actors and Institutions"), 6))
        add_slide("Data", "The next analytical layer is data and indicators.", _bullets(_read_section(markdown, "Data and Indicators"), 6))
        add_slide("Opportunities", "The opportunity set is where strategy becomes concrete.", _bullets(_read_section(markdown, "Opportunities"), 5))
        add_slide("Risks", "The risks show where the story can break.", _bullets(_read_section(markdown, "Risks"), 5))
        add_slide("Watch next", "The follow-up agenda should drive the next research cycle.", _bullets(_read_section(markdown, "Open Questions"), 6))

    output = Path(output_path) if output_path else paths.PRESENTATIONS / f"{topic_slug}-generated-deck.pptx"
    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output)
    return output
